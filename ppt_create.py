import glob
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# 创建一个PPT对象
prs = Presentation()

# 图片文件路径列表
start_ppt = 1
end_ppt = 14
for ppt_no in range(start_ppt, end_ppt+1):
    image_paths = glob.glob(f'./fig/{ppt_no}_*.jpg')
    print(len(image_paths))
    # 遍历图片路径列表，将每张图片插入到PPT中
    for count, image_path in enumerate(image_paths):
        print(f'正在处理第{count+1}/{len(image_paths)}张图片：{image_path}')

        img = Image.open(image_path)
        width, height = img.size
        # 计算调整比例
        max_width = Inches(56)  # 最大幻灯片宽度为56英寸
        max_height = Inches(56)  # 最大幻灯片高度为56英寸
        ratio = min(max_width / width, max_height / height)
        new_width = int(width * ratio)
        new_height = int(height * ratio)

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白幻灯片布局
        # 保持幻灯片和图片比例一致
        prs.slide_width = new_width
        prs.slide_height = new_height
        left = top = Inches(0)
        slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)

    # 保存PPT文件
    prs.save(f'./ppt/{ppt_no}.pptx')
